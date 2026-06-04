from __future__ import annotations

import base64
import html
import json
import os
from pathlib import Path

import pandas as pd
from PIL import Image, ImageChops
from pptx import Presentation
from pptx.util import Inches
from playwright.sync_api import sync_playwright
from diagrams import Cluster, Diagram, Edge
from diagrams.onprem.client import Users
from diagrams.onprem.compute import Server
from diagrams.onprem.inmemory import Redis
from diagrams.onprem.monitoring import Grafana, Prometheus
from diagrams.onprem.queue import Kafka
from diagrams.onprem.analytics import Flink
from diagrams.onprem.database import Clickhouse
from diagrams.onprem.tracing import Jaeger
from diagrams.aws.storage import S3


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "outputs" / "presentations"
CHARTS = ROOT / "outputs" / "charts"
NB_ASSETS = OUT / "notebook-assets"
HTML_OUT = OUT / "shopx-aiops-final-html.html"
PPTX_OUT = OUT / "shopx-aiops-final-rca-html.pptx"
PPTX_FALLBACK = OUT / "shopx-aiops-final-rca-html-new.pptx"
RENDER_DIR = OUT / "html-slide-renders"
PIPELINE_DIAGRAM = CHARTS / "streaming-pipeline-diagram.png"


def b64(path: Path) -> str:
    ext = path.suffix.lower().lstrip(".")
    mime = "jpeg" if ext in {"jpg", "jpeg"} else ext
    return f"data:image/{mime};base64,{base64.b64encode(path.read_bytes()).decode('ascii')}"


def save_notebook_outputs() -> list[Path]:
    NB_ASSETS.mkdir(parents=True, exist_ok=True)
    nb = json.loads((ROOT / "EDA.ipynb").read_text(encoding="utf-8"))
    paths: list[Path] = []
    fig_no = 1
    for cell in nb.get("cells", []):
        for output in cell.get("outputs", []):
            data = output.get("data") or {}
            if "image/png" not in data:
                continue
            path = NB_ASSETS / f"notebook-figure-{fig_no:02d}.png"
            path.write_bytes(base64.b64decode(data["image/png"]))
            paths.append(path)
            fig_no += 1
    return paths


def build_pipeline_diagram() -> Path:
    CHARTS.mkdir(parents=True, exist_ok=True)
    graphviz_bin = Path(r"C:\Program Files\Graphviz\bin")
    if graphviz_bin.exists():
        os.environ["PATH"] = str(graphviz_bin) + os.pathsep + os.environ.get("PATH", "")
    target_no_ext = CHARTS / "streaming-pipeline-diagram"
    with Diagram(
        "ShopX Live Anomaly Detection and RCA Pipeline",
        filename=str(target_no_ext),
        outformat="png",
        show=False,
        direction="LR",
        graph_attr={
            "bgcolor": "transparent",
            "pad": "0.45",
            "nodesep": "1.0",
            "ranksep": "1.1",
            "splines": "spline",
            "fontname": "Arial",
            "fontsize": "24",
            "dpi": "180",
        },
        node_attr={"fontname": "Arial", "fontsize": "18"},
        edge_attr={"fontname": "Arial", "fontsize": "15", "color": "#52606d", "penwidth": "2"},
    ):
        services = Server("ShopX services\napi/cart/order/payment")
        otel = Prometheus("OpenTelemetry\nCollector")
        kafka = Kafka("Kafka\ntelemetry topics")

        with Cluster("Stream processing"):
            flink = Flink("Flink windows\nfeature extraction")
            mad = Redis("MAD + IF\nmetric alerts")
            drain = Redis("Drain3\nlog templates")

        with Cluster("Stores"):
            hot = Clickhouse("Hot stores\nmetrics + logs")
            cold = S3("S3 Parquet\nreplay archive")

        dash = Grafana("RCA dashboard\nWHEN / WHERE / WHAT")

        services >> Edge(label="metrics / logs / traces") >> otel
        otel >> Edge(label="batch + enrich") >> kafka
        kafka >> Edge(label="stream") >> flink
        flink >> mad
        flink >> drain
        kafka >> Edge(label="raw replay") >> cold
        [mad, drain] >> Edge(label="evidence") >> hot
        [hot, cold] >> Edge(label="query + replay") >> dash
    image = Image.open(PIPELINE_DIAGRAM).convert("RGBA")
    alpha_bbox = image.getchannel("A").getbbox()
    rgb = image.convert("RGB")
    white = Image.new("RGB", rgb.size, "white")
    diff = ImageChops.difference(rgb, white)
    white_bbox = diff.point(lambda p: 255 if p > 12 else 0).getbbox()
    bbox = white_bbox or alpha_bbox
    if bbox:
        image = image.crop(bbox)
        pad = 36
        padded = Image.new("RGBA", (image.width + pad * 2, image.height + pad * 2), (255, 255, 255, 0))
        padded.paste(image, (pad, pad), image)
        padded.save(PIPELINE_DIAGRAM)
    return PIPELINE_DIAGRAM


def metric_card_stats() -> dict[str, str]:
    anom = pd.read_csv(ROOT / "outputs" / "anomalies_metrics.csv")
    rows = {}
    keys = {
        "5xx": ("http_5xx_rate", "2026-06-01T06:08:00+00:00"),
        "latency": ("http_p99_latency_ms", "2026-06-01T14:40:00+00:00"),
        "memory": ("memory_usage_bytes", "2026-06-01T16:26:00+00:00"),
        "gc": ("jvm_gc_pause_ms_avg", "2026-06-01T17:50:30+00:00"),
    }
    for name, (metric, timestamp) in keys.items():
        r = anom[(anom.service == "cart-service") & (anom.metric == metric) & (anom.timestamp == timestamp)].iloc[0]
        rows[name] = {
            "time": timestamp[11:19],
            "metric": metric,
            "value": float(r.value),
            "threshold": float(r.threshold),
            "score": None if pd.isna(r.score) else float(r.score),
        }
    return rows


def notebook_text_blocks() -> tuple[str, str]:
    nb = json.loads((ROOT / "EDA.ipynb").read_text(encoding="utf-8"))
    loaded = ""
    stats = ""
    for i, cell in enumerate(nb.get("cells", [])):
        texts = []
        for output in cell.get("outputs", []):
            if "text" in output:
                val = output["text"]
                texts.append("".join(val) if isinstance(val, list) else str(val))
            data = output.get("data") or {}
            if "text/plain" in data:
                val = data["text/plain"]
                texts.append("".join(val) if isinstance(val, list) else str(val))
        if i == 0:
            loaded = "\n".join(texts).strip()
        if i == 1:
            stats = "\n".join(texts).strip()
    return loaded, stats


def table_rows(path: Path, cols: list[str], limit: int | None = None, where=None) -> str:
    df = pd.read_csv(path)
    if where is not None:
        df = where(df)
    if limit:
        df = df.head(limit)
    out = []
    for _, row in df[cols].iterrows():
        out.append("<tr>" + "".join(f"<td>{html.escape(str(row[c]))}</td>" for c in cols) + "</tr>")
    return "\n".join(out)


def build_html() -> str:
    build_pipeline_diagram()
    nb_figs = save_notebook_outputs()
    loaded, stats_text = notebook_text_blocks()
    s = metric_card_stats()
    cart_eda = pd.read_csv(CHARTS / "cart_eda_distribution_stats.csv")
    eda_rows = "\n".join(
        f"<tr><td>{html.escape(r.metric)}</td><td>{r.skewness:.2f}</td><td>{r.kurtosis:.2f}</td><td>{r.p50:.2f}</td><td>{r.p99:.2f}</td></tr>"
        for r in cart_eda.itertuples()
    )
    mad_rows = table_rows(
        ROOT / "outputs" / "anomalies_metrics.csv",
        ["timestamp", "metric", "value", "baseline_median", "threshold", "score"],
        where=lambda df: df[(df.service == "cart-service") & (df.metric.isin(["http_5xx_rate", "http_p99_latency_ms", "memory_usage_bytes", "jvm_gc_pause_ms_avg", "container_restart_count"]))].drop_duplicates("metric"),
    )
    log_rows = table_rows(
        ROOT / "outputs" / "log_templates.csv",
        ["template_id", "service", "level", "count", "first_seen", "template"],
        where=lambda df: df[df.template.str.contains("GC overhead|ProductCatalogCache eviction|OOMKilled|Cart service timeout|Cart service returned 5xx", case=False, na=False)],
    )
    timeline_rows = table_rows(ROOT / "outputs" / "incident_timeline.csv", ["timestamp", "service", "signal", "evidence"], limit=12)

    css = """
    :root{--bg:#080b10;--panel:#111821;--panel2:#172130;--ink:#f5f2ea;--muted:#9aa7b4;--line:#2b3746;--red:#ff5d4a;--orange:#ffb14a;--blue:#59a6ff;--green:#51d19d}
    *{box-sizing:border-box} body{margin:0;background:#030507;color:var(--ink);font-family:"Segoe UI",Arial,sans-serif;text-rendering:geometricPrecision}
    .deck{width:1600px;margin:0 auto}.slide{width:1600px;height:900px;position:relative;overflow:hidden;padding:58px 72px;background:radial-gradient(circle at 82% 12%,#26384d 0,#101722 26%,#080b10 58%);}
    .slide.light{background:linear-gradient(135deg,#f7f0e3,#fefcf7);color:#13202c}.slide.light .eyebrow{color:#8e3b2f}.slide.light .sub,.slide.light .muted{color:#52606d}.slide.light .card,.slide.light .panel{background:#fffdf8;border-color:#dfd4c4;color:#13202c}.slide.light table td,.slide.light table th{border-color:#dfd4c4}
    .eyebrow{font-size:18px;line-height:1.35;letter-spacing:.08em;text-transform:uppercase;color:#ff725f;font-weight:800}.title{font-size:52px;line-height:1.16;font-weight:800;letter-spacing:0;max-width:1140px;margin:24px 0 0}.sub{font-size:22px;line-height:1.45;color:var(--muted);max-width:1040px;margin-top:18px}
    .grid{display:grid;gap:22px}.cols3{grid-template-columns:repeat(3,1fr)}.cols4{grid-template-columns:repeat(4,1fr)}.card,.panel{background:linear-gradient(180deg,rgba(255,255,255,.075),rgba(255,255,255,.035));border:1px solid var(--line);border-radius:22px;padding:24px;box-shadow:0 20px 60px rgba(0,0,0,.18)}.card h3{margin:0 0 14px;font-size:18px;line-height:1.28;color:var(--muted);text-transform:uppercase;letter-spacing:.06em}.big{font-size:42px;line-height:1.16;font-weight:800}.note{font-size:20px;line-height:1.42;color:var(--muted)}.accent{color:var(--red)}.blue{color:var(--blue)}.green{color:var(--green)}.orange{color:var(--orange)}
    .footer{position:absolute;left:72px;right:72px;bottom:30px;border-top:1px solid rgba(150,160,170,.25);padding-top:12px;color:#7f8b98;font-size:13px;display:flex;justify-content:space-between}
    .figure{background:white;border-radius:18px;padding:14px;border:1px solid rgba(255,255,255,.2)}.figure img{width:100%;height:100%;object-fit:contain;display:block}.code{white-space:pre-wrap;font-family:Consolas,monospace;font-size:18px;line-height:1.35;background:#0c121a;border:1px solid #263342;border-radius:18px;padding:22px;color:#dbe7f3}
    table{width:100%;border-collapse:collapse;font-size:16px}th,td{border-bottom:1px solid var(--line);padding:10px 12px;text-align:left;vertical-align:top}th{font-size:13px;color:var(--muted);text-transform:uppercase;letter-spacing:.05em}.timeline{position:absolute;left:130px;top:240px;width:1220px}.step{display:grid;grid-template-columns:130px 170px 1fr;gap:22px;align-items:center;margin:19px 0}.dot{width:22px;height:22px;border-radius:50%;background:var(--red);box-shadow:0 0 0 8px rgba(255,93,74,.12)}.rail{position:absolute;left:249px;top:258px;width:4px;height:470px;background:#344252}.tag{font-size:15px;line-height:1.3;color:var(--muted);text-transform:uppercase;font-weight:800}.time{font-size:24px;line-height:1.25;font-weight:800}.desc{font-size:23px;line-height:1.32}.kpi{font-size:76px;line-height:1.08;font-weight:800;letter-spacing:0}.split{display:grid;grid-template-columns:1.1fr .9fr;gap:34px;align-items:center}.small{font-size:15px}.quote{font-size:32px;line-height:1.28;font-weight:800}.pill{display:inline-block;border:1px solid var(--line);border-radius:999px;padding:10px 15px;color:var(--muted);font-weight:800;margin-right:8px}.pipeline{display:grid;grid-template-columns:repeat(5,1fr);gap:18px;margin-top:62px}.pipe-card{height:210px}.arrow{position:absolute;top:390px;width:34px;height:4px;background:#617085}.arrow:after{content:"";position:absolute;right:-8px;top:-6px;border-left:10px solid #617085;border-top:8px solid transparent;border-bottom:8px solid transparent}
    .flow{display:grid;gap:18px;margin-top:38px}.flow5{grid-template-columns:repeat(5,1fr)}.flow4{grid-template-columns:repeat(4,1fr)}.flow3{grid-template-columns:repeat(3,1fr)}.flow-card{position:relative;min-height:168px}.flow-card:after{content:"";position:absolute;right:-18px;top:75px;width:18px;height:3px;background:#9aa7b4}.flow-card:last-child:after{display:none}.num{width:42px;height:42px;border-radius:14px;background:var(--red);color:white;display:grid;place-items:center;font-weight:900;font-size:22px;margin-bottom:18px}.flow-card h3{font-size:22px;color:var(--ink);text-transform:none;letter-spacing:0}.dark .flow-card h3{color:var(--ink)}.live-grid{display:grid;grid-template-columns:1fr 1fr;gap:24px;margin-top:34px}.lane-title{font-size:24px;font-weight:900;margin-bottom:12px}.mini-list{font-size:19px;line-height:1.38;color:var(--muted);margin:0;padding-left:22px}.matrix{display:grid;grid-template-columns:260px 1fr 1fr;gap:0;margin-top:34px;border:1px solid var(--line);border-radius:22px;overflow:hidden}.matrix>div{padding:18px 20px;border-right:1px solid var(--line);border-bottom:1px solid var(--line);font-size:18px}.matrix>div:nth-child(3n){border-right:0}.matrix .head{font-weight:900;background:rgba(255,255,255,.08);color:var(--ink)}.light .matrix .head{background:#efe7da}.matrix .label{font-weight:900;color:var(--red)}
    """

    slides = []
    def footer(n, source):
        return f'<div class="footer"><span>{html.escape(source)}</span><span>{n:02d}</span></div>'

    slides.append(f"""
    <section class="slide">
      <div class="eyebrow">ShopX AIOps W1 RCA</div>
      <div class="title">cart-service không chỉ timeout; nó đi vào heap/cache pressure trước khi OOM.</div>
      <div class="sub">Evidence-first final deck: EDA shape, detector threshold, Drain3 log templates, notebook output và production pipeline.</div>
      <div class="grid cols3" style="position:absolute;left:72px;right:72px;bottom:145px">
        <div class="card"><h3>WHEN</h3><div class="big accent">06:08 UTC</div><div class="note">cart 5xx value 1.03 &gt; MAD threshold 0.354</div></div>
        <div class="card"><h3>WHERE</h3><div class="big blue">cart-service</div><div class="note">GC/cache/OOM templates xuất hiện trước downstream</div></div>
        <div class="card"><h3>WHAT</h3><div class="big green">heap/cache pressure</div><div class="note">ProductCatalogCache eviction failure + OOMKilled + restart loop</div></div>
      </div>{footer(1,'FINDINGS.md, anomalies_metrics.csv, incident_timeline.csv')}
    </section>""")

    slides.append(f"""
    <section class="slide light">
      <div class="eyebrow">Notebook output: data loaded</div>
      <div class="title" style="font-size:50px">EDA notebook xác nhận dữ liệu đủ sạch để phân tích, không phải dữ liệu log/metric rác.</div>
      <div class="split" style="margin-top:44px">
        <div class="code" style="height:390px">{html.escape(loaded)}</div>
        <div class="grid" style="gap:18px">
          <div class="card"><h3>Metric coverage</h3><div class="kpi accent">2,820</div><div class="note">rows / service, 30s interval, one full day</div></div>
          <div class="card"><h3>Missing values</h3><div class="kpi green">0%</div><div class="note">notebook missing table: all key metrics have missing_count = 0</div></div>
        </div>
      </div>{footer(2,'EDA.ipynb cell outputs')}
    </section>""")

    slides.append(f"""
    <section class="slide light">
      <div class="eyebrow">EDA: distribution shape</div>
      <div class="title" style="font-size:48px">Right-skew không phải nói chung chung: chính dataset này có tail phải rất rõ.</div>
      <div class="split" style="grid-template-columns:.8fr 1.2fr;margin-top:42px">
        <div class="panel">
          <table><thead><tr><th>cart metric</th><th>skew</th><th>kurtosis</th><th>p50</th><th>p99</th></tr></thead><tbody>{eda_rows}</tbody></table>
        </div>
        <div class="figure" style="height:440px"><img src="{b64(CHARTS/'eda-cart-distributions.png')}"></div>
      </div>
      <div class="sub" style="font-size:22px;max-width:1320px">Kurtosis cao = nhiều extreme values/outlier hơn normal shape. Vì vậy median/MAD ít bị kéo lệch hơn mean/std.</div>
      {footer(3,'outputs/charts/eda-cart-distributions.png, cart_eda_distribution_stats.csv')}
    </section>""")

    slides.append(f"""
    <section class="slide">
      <div class="eyebrow">Notebook figures embedded</div>
      <div class="title" style="font-size:46px">Notebook output được đưa thẳng vào deck: histogram/boxplot cho các signal chính.</div>
      <div class="grid cols4" style="margin-top:42px">
        {''.join(f'<div class="figure" style="height:260px"><img src="{b64(p)}"></div>' for p in nb_figs[:4])}
      </div>
      <div class="sub">Các figure này là output trực tiếp từ `EDA.ipynb`, dùng để support phân phối skewed và outlier.</div>
      {footer(4,'EDA.ipynb image/png outputs')}
    </section>""")

    slides.append(f"""
    <section class="slide light">
      <div class="eyebrow">Exact robust MAD evidence</div>
      <div class="title" style="font-size:46px">Timestamp anomaly map trực tiếp vào value, median, threshold và score.</div>
      <div class="figure" style="height:420px;margin-top:42px"><img src="{b64(CHARTS/'mad-cart-exact-evidence.png')}"></div>
      <div class="grid cols4" style="margin-top:22px">
        <div class="card"><h3>5xx</h3><div class="big accent">{s['5xx']['time']}</div><div class="note">1.03 &gt; 0.354</div></div>
        <div class="card"><h3>Latency</h3><div class="big orange">{s['latency']['time']}</div><div class="note">148.7ms &gt; 122.8ms</div></div>
        <div class="card"><h3>Memory</h3><div class="big blue">{s['memory']['time']}</div><div class="note">0.62GB &gt; 0.57GB</div></div>
        <div class="card"><h3>GC pause</h3><div class="big green">{s['gc']['time']}</div><div class="note">131.8ms &gt; 104.3ms</div></div>
      </div>{footer(5,'outputs/anomalies_metrics.csv')}
    </section>""")

    slides.append(f"""
    <section class="slide">
      <div class="eyebrow">Detector decision</div>
      <div class="title" style="font-size:48px">Final stance: robust MAD + IsolationForest; EWMA chỉ dùng để đọc drift.</div>
      <div class="split" style="margin-top:40px">
        <div class="figure" style="height:470px"><img src="{b64(CHARTS/'detector-evidence-table.png')}"></div>
        <div class="grid">
          <div class="card"><h3>MAD</h3><div class="big accent">primary</div><div class="note">threshold cụ thể theo metric, audit được WHEN</div></div>
          <div class="card"><h3>IsolationForest</h3><div class="big blue">07:27</div><div class="note">cart-service multivariate abnormality confirmation</div></div>
          <div class="card"><h3>EWMA caveat</h3><div class="big orange">false positives</div><div class="note">payment timeout 09:51, cart 5xx baseline FP=12; không dùng làm RCA chính</div></div>
        </div>
      </div>{footer(6,'method_comparison.csv, detector_observability.csv')}
    </section>""")

    slides.append(f"""
    <section class="slide light">
      <div class="eyebrow">Metric evidence</div>
      <div class="title" style="font-size:48px">Cart metrics cho thấy pressure tích lũy trước OOM/restart.</div>
      <div class="split" style="grid-template-columns:1.25fr .75fr;margin-top:34px">
        <div class="figure" style="height:515px"><img src="{b64(CHARTS/'cart-service-metrics.png')}"></div>
        <div class="grid">
          <div class="card"><h3>Memory pressure</h3><div class="big blue">16:26</div><div class="note">memory_usage_bytes vượt MAD threshold</div></div>
          <div class="card"><h3>GC pressure</h3><div class="big orange">17:50</div><div class="note">GC pause 131.8ms &gt; threshold 104.3ms</div></div>
          <div class="card"><h3>Visible failure</h3><div class="big accent">20:00</div><div class="note">restart_count từ 0 lên 1 sau OOMKilled</div></div>
        </div>
      </div>{footer(7,'cart-service-metrics.png, anomalies_metrics.csv')}
    </section>""")

    slides.append(f"""
    <section class="slide">
      <div class="eyebrow">Log evidence + Drain3 fit</div>
      <div class="title" style="font-size:46px">Logs structured và consistent; Drain3 dùng để template hóa dynamic params và count pattern theo thời gian.</div>
      <div class="split" style="margin-top:34px">
        <div class="figure" style="height:455px"><img src="{b64(CHARTS/'key-log-template-timeseries.png')}"></div>
        <div class="panel" style="height:455px;overflow:hidden">
          <table><thead><tr><th>ID</th><th>svc</th><th>lvl</th><th>count</th><th>first seen</th><th>template</th></tr></thead><tbody>{log_rows}</tbody></table>
        </div>
      </div>
      <div class="sub" style="font-size:22px">Bạn nói đúng: logs khá gọn. Drain3 không “cứu data bẩn”; nó biến message có userId/status/duration/heap/pause thành template đếm được.</div>
      {footer(8,'g1/logs/*.jsonl, log_templates.csv, key-log-template-timeseries.png')}
    </section>""")

    steps = [
        ("06:08", "cart metric", "http_5xx_rate anomaly: value 1.03 > threshold 0.354"),
        ("06:30", "cart log", "GC overhead warning, pause=713ms heap=93%"),
        ("06:33", "cart log", "ProductCatalogCache eviction failed: heap pressure too high"),
        ("16:26", "cart metric", "memory_usage_bytes crosses robust threshold"),
        ("17:50", "cart metric", "jvm_gc_pause_ms_avg crosses threshold"),
        ("19:59", "cart log", "OOMKilled: memory limit exceeded"),
        ("20:00+", "fan-out", "restart loop, gateway upstream error, order/payment timeout"),
    ]
    steps_html = '<div class="rail"></div>' + "".join(f'<div class="step"><div class="time">{t}</div><div><span class="dot"></span></div><div><div class="tag">{html.escape(tag)}</div><div class="desc">{html.escape(desc)}</div></div></div>' for t, tag, desc in steps)
    slides.append(f"""
    <section class="slide">
      <div class="eyebrow">Evidence ordering</div>
      <div class="title" style="font-size:48px">RCA mạnh lên nhờ thứ tự bằng chứng, không nhờ một signal đơn lẻ.</div>
      <div class="timeline">{steps_html}</div>
      {footer(9,'incident_timeline.csv, FINDINGS.md')}
    </section>""")

    slides.append(f"""
    <section class="slide light">
      <div class="eyebrow">Python diagrams asset</div>
      <div class="title" style="font-size:42px;max-width:1320px">Generated production pipeline diagram, simplified for slide readability.</div>
      <div class="split" style="grid-template-columns:1.62fr .58fr;margin-top:24px;align-items:stretch">
        <div class="figure" style="height:470px;padding:16px"><img src="{b64(PIPELINE_DIAGRAM)}"></div>
        <div class="grid" style="gap:14px">
          <div class="card" style="padding:18px 20px;min-height:140px"><h3>Ingest</h3><div class="big blue" style="font-size:30px">OTel -> Kafka</div><div class="note" style="font-size:17px">Collector enriches; Kafka buffers and enables replay.</div></div>
          <div class="card" style="padding:18px 20px;min-height:140px"><h3>Detect</h3><div class="big accent" style="font-size:30px">Flink + MAD/IF</div><div class="note" style="font-size:17px">Windowed features, metric alerts, and Drain3 templates.</div></div>
          <div class="card" style="padding:18px 20px;min-height:140px"><h3>RCA</h3><div class="big green" style="font-size:30px">Dashboard</div><div class="note" style="font-size:17px">Hot stores and replay archive support evidence ordering.</div></div>
        </div>
      </div>
      {footer(10,'Generated by Python diagrams package in tools/build_html_powerpoint.py')}
    </section>""")

    slides.append(f"""
    <section class="slide light">
      <div class="eyebrow">Live anomaly detection + RCA pipeline</div>
      <div class="title" style="font-size:44px;max-width:1280px">Production flow: telemetry được ingest liên tục, detector chạy theo window, RCA ghép evidence theo thời gian.</div>
      <div class="flow flow5">
        <div class="card flow-card"><div class="num">1</div><h3>Emit</h3><div class="note">Services emit metrics, logs, traces with service/pod labels.</div></div>
        <div class="card flow-card"><div class="num">2</div><h3>Ingest</h3><div class="note">OpenTelemetry Collector batches, enriches, filters, and routes signals.</div></div>
        <div class="card flow-card"><div class="num">3</div><h3>Buffer</h3><div class="note">Kafka topics decouple services from processing and enable replay.</div></div>
        <div class="card flow-card"><div class="num">4</div><h3>Detect</h3><div class="note">Flink computes windows; MAD + IF score metrics; Drain3 templates logs.</div></div>
        <div class="card flow-card"><div class="num">5</div><h3>RCA</h3><div class="note">Evidence timeline ranks origin candidate and downstream blast radius.</div></div>
      </div>
      <div class="panel" style="margin-top:42px">
        <span class="quote" style="font-size:30px">Early alert target: memory slope + GC pause + cache eviction template count + cart p99/5xx.</span>
      </div>
      {footer(11,'Production scenario synthesized from ARCHITECTURE.md and DATA_PIPELINE_PRESENTATION.md')}
    </section>""")

    slides.append(f"""
    <section class="slide">
      <div class="eyebrow">What we currently have in this repo</div>
      <div class="title" style="font-size:44px;max-width:1320px">Our live pipeline is a local replay simulator: same RCA logic, but file-backed instead of Kafka/Flink-backed.</div>
      <div class="live-grid">
        <div class="panel">
          <div class="lane-title blue">Current implementation</div>
          <ul class="mini-list">
            <li><b>Source:</b> `g1/metrics/*.csv` + `g1/logs/*.jsonl`.</li>
            <li><b>Streamer:</b> `w1/lab/realtime.py` replays rows in timestamp order.</li>
            <li><b>Signals:</b> rolling metric state + Drain3 template events.</li>
            <li><b>Outputs:</b> `events.jsonl`, `alerts.jsonl`, `signals.json`, RCA JSON files.</li>
          </ul>
        </div>
        <div class="panel">
          <div class="lane-title green">What production would swap in</div>
          <ul class="mini-list">
            <li><b>Source:</b> services instrumented with OTel SDK.</li>
            <li><b>Transport:</b> Kafka topics instead of local files.</li>
            <li><b>Processing:</b> Flink jobs instead of Python replay loop.</li>
            <li><b>Serving:</b> dashboard reads hot stores and RCA state continuously.</li>
          </ul>
        </div>
      </div>
      <div class="matrix">
        <div class="head">Pipeline concern</div><div class="head">Current repo</div><div class="head">Production scenario</div>
        <div class="label">Ingest</div><div>Local CSV/JSONL loaders and replay events</div><div>OTel Collector receives live telemetry</div>
        <div class="label">Streaming</div><div>`stream_events.jsonl` simulates event flow</div><div>Kafka partitions by signal/service</div>
        <div class="label">Processing</div><div>Python state machine, MAD thresholds, Drain3 templates</div><div>Flink window jobs + model workers</div>
        <div class="label">RCA</div><div>Rule-based evidence gate writes `rca_timeline.json`</div><div>RCA service correlates metric/log/trace evidence continuously</div>
      </div>
      {footer(12,'w1/lab/realtime.py, outputs/realtime/*.jsonl, outputs/realtime/*.json')}
    </section>""")

    slides.append(f"""
    <section class="slide light">
      <div class="eyebrow">Simulating the data flow</div>
      <div class="title" style="font-size:44px;max-width:1320px">The demo button replays the incident, emits alerts, then updates the RCA hypothesis.</div>
      <div class="flow flow4">
        <div class="card flow-card"><div class="num">A</div><h3>Replay</h3><div class="note">Read metric rows and Drain3 log-template events in timestamp order.</div></div>
        <div class="card flow-card"><div class="num">B</div><h3>Gate</h3><div class="note">Apply baseline thresholds, 3-of-5 persistence, and log-template count gates.</div></div>
        <div class="card flow-card"><div class="num">C</div><h3>Correlate</h3><div class="note">Merge cart memory/GC/cache/OOM signals with downstream timeout/5xx.</div></div>
        <div class="card flow-card"><div class="num">D</div><h3>Present</h3><div class="note">Dashboard displays alert stream, RCA timeline, and ranked hypothesis.</div></div>
      </div>
      <div class="grid cols3" style="margin-top:44px">
        <div class="card"><h3>Data stream</h3><div class="big blue">events.jsonl</div><div class="note">all simulated telemetry events</div></div>
        <div class="card"><h3>Detection stream</h3><div class="big accent">alerts.jsonl</div><div class="note">metric + log-template alerts</div></div>
        <div class="card"><h3>RCA output</h3><div class="big green">hypotheses.json</div><div class="note">cart-service ranked as origin candidate</div></div>
      </div>
      {footer(13,'realtime_dashboard.py snapshot and realtime.py replay pipeline')}
    </section>""")

    slides.append(f"""
    <section class="slide light">
      <div class="eyebrow">ADR-style decisions</div>
      <div class="title" style="font-size:46px">Architecture decisions: đủ Context / Decision / Rationale / Trade-off.</div>
      <div class="grid cols3" style="margin-top:38px">
        <div class="card"><h3>ADR-01: Collection</h3><div class="big blue">OTel</div><div class="note"><b>Context:</b> nhiều service, cần schema chung.<br><b>Decision:</b> SDK + Collector.<br><b>Trade-off:</b> thêm instrumentation effort.</div></div>
        <div class="card"><h3>ADR-02: Transport</h3><div class="big orange">Kafka</div><div class="note"><b>Context:</b> telemetry burst/backpressure.<br><b>Decision:</b> topics by signal type.<br><b>Trade-off:</b> operational complexity.</div></div>
        <div class="card"><h3>ADR-03: Processing</h3><div class="big accent">Flink</div><div class="note"><b>Context:</b> cần rolling windows/replay.<br><b>Decision:</b> streaming; pandas for lab.<br><b>Trade-off:</b> heavier than batch.</div></div>
        <div class="card"><h3>ADR-04: Detection</h3><div class="big green">MAD + IF</div><div class="note"><b>Context:</b> skewed incident metrics.<br><b>Decision:</b> MAD primary, IF confirm.<br><b>Trade-off:</b> IF less explainable.</div></div>
        <div class="card"><h3>ADR-05: Logs</h3><div class="big blue">Drain3</div><div class="note"><b>Context:</b> logs structured but dynamic params.<br><b>Decision:</b> template mining.<br><b>Trade-off:</b> calibration needed.</div></div>
        <div class="card"><h3>ADR-06: Storage</h3><div class="big orange">Hot + cold</div><div class="note"><b>Context:</b> query recent + replay history.<br><b>Decision:</b> VM/Loki/S3.<br><b>Trade-off:</b> multiple stores.</div></div>
      </div>
      {footer(14,'ADR content synthesized from ARCHITECTURE.md and DATA_PIPELINE_PRESENTATION.md')}
    </section>""")

    slides.append(f"""
    <section class="slide">
      <div class="eyebrow">Final takeaway</div>
      <div class="title">cart-service là origin candidate; prevention phải bắt heap/cache pressure trước OOM.</div>
      <div class="grid cols3" style="margin-top:78px">
        <div class="card"><h3>RCA</h3><div class="big accent">heap/cache</div><div class="note">GC warning + cache eviction failure + OOMKilled</div></div>
        <div class="card"><h3>Detection</h3><div class="big blue">MAD + IF</div><div class="note">MAD cho threshold; IF confirm service abnormality</div></div>
        <div class="card"><h3>Ops</h3><div class="big green">guardrail</div><div class="note">restart loop + OOMKilled + downstream blast radius</div></div>
      </div>
      <div class="panel" style="position:absolute;left:120px;right:120px;bottom:140px;text-align:center"><span class="quote">Metrics answer WHEN. Log templates answer WHERE. Evidence ordering supports WHAT.</span></div>
      {footer(15,'Final synthesis from repo artifacts')}
    </section>""")

    return f"<!doctype html><html><head><meta charset='utf-8'><style>{css}</style></head><body><main class='deck'>{''.join(slides)}</main></body></html>"


def render_and_export() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    RENDER_DIR.mkdir(parents=True, exist_ok=True)
    HTML_OUT.write_text(build_html(), encoding="utf-8")
    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={"width": 1600, "height": 900}, device_scale_factor=1)
        page.goto(HTML_OUT.resolve().as_uri())
        count = page.locator(".slide").count()
        pngs = []
        for i in range(count):
            path = RENDER_DIR / f"slide-{i+1:02d}.png"
            page.locator(".slide").nth(i).screenshot(path=str(path))
            pngs.append(path)
        browser.close()

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for png in pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(png), 0, 0, width=prs.slide_width, height=prs.slide_height)
    try:
        prs.save(PPTX_OUT)
        final_pptx = PPTX_OUT
    except PermissionError:
        prs.save(PPTX_FALLBACK)
        final_pptx = PPTX_FALLBACK
        print(f"Could not overwrite locked file: {PPTX_OUT}")
    print(f"Wrote {HTML_OUT}")
    print(f"Wrote {final_pptx}")
    print(f"Rendered {len(pngs)} slides")


if __name__ == "__main__":
    render_and_export()
